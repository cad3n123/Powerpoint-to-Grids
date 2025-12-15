import sys
import os
import platform
import subprocess
import time

# We import PyQt5 immediately for the UI, but delay heavy imports (pptx, pdf2image)
# until the thread starts to ensure the app opens instantly.
try:
    from PyQt5.QtWidgets import (QApplication, QWidget, QVBoxLayout, QPushButton, 
                                 QLabel, QProgressBar, QFileDialog, QMessageBox)
    from PyQt5.QtCore import QThread, pyqtSignal, Qt
except ImportError:
    pass # Builder will handle dependencies

# --- WORKER THREAD (Runs in background) ---
class Worker(QThread):
    progress = pyqtSignal(int)
    status = pyqtSignal(str)
    finished = pyqtSignal()
    error = pyqtSignal(str)

    def __init__(self, files, output_dir):
        super().__init__()
        self.files = files
        self.output_dir = output_dir
        self.is_running = True

    def run(self):
        try:
            # --- DELAYED IMPORTS (Speed up startup) ---
            self.status.emit("Loading libraries...")
            from pptx import Presentation
            from pptx.util import Inches
            from pdf2image import convert_from_path
            
            # --- CONFIG ---
            ROWS = 4
            COLS = 5
            ITEMS_PER_SLIDE = ROWS * COLS
            MARGIN_X = Inches(0.5)
            MARGIN_Y = Inches(0.5)
            SPACING = Inches(0.2)

            # --- HELPERS ---
            def get_poppler_path():
                if getattr(sys, 'frozen', False):
                    return os.path.join(sys._MEIPASS, 'poppler_bin')
                return None

            def get_libreoffice():
                if platform.system() == "Windows":
                    candidates = [
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
                    ]
                    for p in candidates:
                        if os.path.exists(p): return p
                    if subprocess.call("where soffice", shell=True, stdout=subprocess.DEVNULL) == 0:
                        return "soffice"
                return None

            def run_powershell(cmd):
                try:
                    full_cmd = ["powershell", "-NoProfile", "-NonInteractive", "-Command", cmd]
                    subprocess.run(full_cmd, capture_output=True, creationflags=subprocess.CREATE_NO_WINDOW)
                except: pass

            # --- PRINTER SAFETY ---
            original_printer = None
            safe_printer = "Microsoft Print to PDF"
            
            if platform.system() == "Windows":
                self.status.emit("Checking printer settings...")
                # Get current default
                cmd_get = "Get-CimInstance Win32_Printer | Where-Object Default | Select-Object -ExpandProperty Name"
                try:
                    res = subprocess.run(["powershell", "-NoProfile", "-Command", cmd_get], 
                                       capture_output=True, text=True, creationflags=subprocess.CREATE_NO_WINDOW)
                    original_printer = res.stdout.strip()
                except: pass

                # Switch if needed
                if original_printer and original_printer != safe_printer:
                    self.status.emit("Switching to safe printer...")
                    run_powershell(f"(New-Object -ComObject WScript.Network).SetDefaultPrinter('{safe_printer}')")

            # --- PROCESSING ---
            libreoffice = get_libreoffice()
            if not libreoffice:
                self.error.emit("LibreOffice not found.")
                return

            poppler = get_poppler_path()
            all_images = []
            
            total_steps = len(self.files) * 2 # PDF + Convert
            current_step = 0

            if not os.path.exists(self.output_dir):
                os.makedirs(self.output_dir)

            for pptx_path in self.files:
                if not self.is_running: break
                
                base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                self.status.emit(f"Converting: {base_name}")
                
                # 1. PPTX -> PDF
                cmd = [libreoffice, "--headless", "--convert-to", "pdf", "--outdir", self.output_dir, pptx_path]
                subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
                
                pdf_path = os.path.join(self.output_dir, f"{base_name}.pdf")
                current_step += 1
                self.progress.emit(int((current_step / total_steps) * 80))

                # 2. PDF -> Images
                if os.path.exists(pdf_path):
                    self.status.emit(f"Extracting slides: {base_name}")
                    pages = convert_from_path(pdf_path, dpi=200, poppler_path=poppler)
                    
                    slide_folder = os.path.join(self.output_dir, base_name)
                    if not os.path.exists(slide_folder): os.makedirs(slide_folder)

                    for i, page in enumerate(pages):
                        num = i + 1
                        page.save(os.path.join(slide_folder, f"Slide_{num}.pdf"), "PDF")
                        img_path = os.path.join(slide_folder, f"Slide_{num}.png")
                        page.save(img_path, "PNG")
                        all_images.append(img_path)
                    
                    try: os.remove(pdf_path)
                    except: pass
                
                current_step += 1
                self.progress.emit(int((current_step / total_steps) * 80))

            # --- GRID CREATION ---
            if all_images and self.is_running:
                self.status.emit("Generating Grid Presentation...")
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                blank_layout = prs.slide_layouts[6]

                cell_w = (prs.slide_width - (2 * MARGIN_X) - ((COLS - 1) * SPACING)) / COLS
                cell_h = cell_w * (9/16)

                slide = None
                for i, img in enumerate(all_images):
                    if i % ITEMS_PER_SLIDE == 0:
                        slide = prs.slides.add_slide(blank_layout)
                    
                    row = (i % ITEMS_PER_SLIDE) // COLS
                    col = (i % ITEMS_PER_SLIDE) % COLS
                    left = MARGIN_X + (col * (cell_w + SPACING))
                    top = MARGIN_Y + (row * (cell_h + SPACING))
                    try: slide.shapes.add_picture(img, left, top, width=cell_w, height=cell_h)
                    except: pass

                out_pptx = os.path.join(self.output_dir, "Grid_Summary.pptx")
                prs.save(out_pptx)
                self.progress.emit(100)

            # --- RESTORE PRINTER ---
            if original_printer:
                self.status.emit("Restoring printer settings...")
                run_powershell(f"(New-Object -ComObject WScript.Network).SetDefaultPrinter('{original_printer}')")

            self.finished.emit()

        except Exception as e:
            self.error.emit(str(e))

    def stop(self):
        self.is_running = False

# --- MAIN UI ---
class GridMakerUI(QWidget):
    def __init__(self):
        super().__init__()
        self.init_ui()
        self.worker = None

    def init_ui(self):
        self.setWindowTitle("Slide Grid Maker")
        self.setGeometry(300, 300, 400, 250)
        
        layout = QVBoxLayout()
        layout.setSpacing(15)
        layout.setContentsMargins(30, 30, 30, 30)

        # Title
        self.label = QLabel("PowerPoint to Grid Converter")
        self.label.setAlignment(Qt.AlignCenter)
        self.label.setStyleSheet("font-size: 16px; font-weight: bold; margin-bottom: 10px;")
        layout.addWidget(self.label)

        # Info Text
        self.info_label = QLabel("Select PPTX files to create a contact sheet.")
        self.info_label.setAlignment(Qt.AlignCenter)
        self.info_label.setStyleSheet("color: #666;")
        layout.addWidget(self.info_label)

        # Progress Bar (Hidden initially)
        self.pbar = QProgressBar()
        self.pbar.setVisible(False)
        layout.addWidget(self.pbar)

        # Button
        self.btn = QPushButton("Select Files & Start")
        self.btn.setCursor(Qt.PointingHandCursor)
        self.btn.setStyleSheet("""
            QPushButton {
                background-color: #0078D7; 
                color: white; 
                padding: 10px; 
                border-radius: 5px;
                font-size: 14px;
            }
            QPushButton:hover { background-color: #005A9E; }
        """)
        self.btn.clicked.connect(self.start_process)
        layout.addWidget(self.btn)

        self.setLayout(layout)

    def start_process(self):
        # 1. Select Files
        files, _ = QFileDialog.getOpenFileNames(self, "Select PowerPoint Files", "", "PowerPoint Files (*.pptx)")
        if not files: return

        # 2. Select Output Directory (Smart Default)
        last_dir = os.path.dirname(files[0])
        default_name = os.path.join(last_dir, "Processed Slides")
        
        # We use getSaveFileName to let user type a folder name easily, 
        # though we will treat it as a directory.
        output_dir, _ = QFileDialog.getSaveFileName(self, "Create Output Folder", default_name, "")
        
        if not output_dir: return

        # Prepare UI for processing
        self.btn.setEnabled(False)
        self.btn.setText("Processing...")
        self.pbar.setVisible(True)
        self.pbar.setValue(0)
        self.info_label.setText("Starting engine...")

        # Start Thread
        self.worker = Worker(files, output_dir)
        self.worker.progress.connect(self.update_progress)
        self.worker.status.connect(self.update_status)
        self.worker.finished.connect(self.process_finished)
        self.worker.error.connect(self.process_error)
        self.worker.start()

    def update_progress(self, val):
        self.pbar.setValue(val)

    def update_status(self, text):
        self.info_label.setText(text)

    def process_finished(self):
        self.pbar.setValue(100)
        self.info_label.setText("Done!")
        self.btn.setText("Select Files & Start")
        self.btn.setEnabled(True)
        QMessageBox.information(self, "Success", "Grid creation complete!")

    def process_error(self, err):
        self.info_label.setText("Error occurred.")
        self.btn.setEnabled(True)
        self.btn.setText("Retry")
        QMessageBox.critical(self, "Error", f"An error occurred:\n{err}")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = GridMakerUI()
    window.show()
    sys.exit(app.exec_())